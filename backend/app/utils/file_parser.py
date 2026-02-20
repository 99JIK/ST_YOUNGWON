from __future__ import annotations

import csv
import io
import json
import logging
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from backend.app.utils.korean_utils import clean_extracted_text

logger = logging.getLogger(__name__)

# 텍스트 추출 가능한 확장자 목록
EXTRACTABLE_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".hwpx",
    ".txt", ".md", ".csv", ".json",
    ".log", ".xml", ".html", ".htm",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".sql", ".sh", ".bat", ".ps1",
}

# 바이너리(텍스트 추출 불가) — 저장만 하고 인덱싱 안 함
BINARY_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg", ".ico", ".webp",
    ".mp3", ".mp4", ".avi", ".mov", ".wav", ".flac",
    ".zip", ".tar", ".gz", ".7z", ".rar",
    ".exe", ".dll", ".so", ".bin",
    ".hwp",  # 구형 HWP(OLE 바이너리) — 별도 라이브러리 필요
}


def is_extractable(filename: str) -> bool:
    """텍스트 추출이 가능한 파일인지 판별합니다."""
    ext = _ext(filename)
    if ext in EXTRACTABLE_EXTENSIONS:
        return True
    # 알 수 없는 확장자면 바이너리 목록에 없으면 텍스트 시도
    if ext not in BINARY_EXTENSIONS:
        return True
    return False


def extract_text(file_path: Path, file_content: bytes | None = None) -> str:
    """파일에서 텍스트를 추출합니다.

    Args:
        file_path: 파일 경로 (확장자 판별 + PDF 등 디스크 필요 시 사용)
        file_content: 파일 바이트 (메모리에 이미 있을 경우)

    Returns:
        추출된 전체 텍스트
    """
    ext = _ext(file_path.name)

    if file_content is None:
        file_content = file_path.read_bytes()

    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext == ".docx":
            return _extract_docx(file_content)
        elif ext == ".xlsx":
            return _extract_xlsx(file_content)
        elif ext == ".pptx":
            return _extract_pptx(file_content)
        elif ext == ".hwpx":
            return _extract_hwpx(file_content)
        elif ext == ".csv":
            return _extract_csv(file_content)
        elif ext == ".json":
            return _extract_json(file_content)
        elif ext in (".xml", ".html", ".htm"):
            return _extract_xml_html(file_content)
        else:
            # 플레인 텍스트 계열 (txt, md, log, py, js, yaml 등)
            return _extract_plain_text(file_content)
    except Exception as e:
        logger.error(f"텍스트 추출 실패 ({file_path.name}): {e}")
        raise


# ─── PDF ───────────────────────────────────────────────
def _extract_pdf(file_path: Path) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(str(file_path))
    pages = []
    for page in doc:
        text = page.get_text("text")
        text = clean_extracted_text(text)
        if text.strip():
            pages.append(text)
    doc.close()
    logger.info(f"PDF 추출 완료: {file_path.name}, {len(pages)}페이지")
    return "\n\n".join(pages)


# ─── DOCX ──────────────────────────────────────────────
def _extract_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # 표(table) 내용도 추출
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    return clean_extracted_text("\n".join(paragraphs))


# ─── XLSX ──────────────────────────────────────────────
def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"[시트: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append(" | ".join(cells))
    wb.close()
    return clean_extracted_text("\n".join(lines))


# ─── PPTX ──────────────────────────────────────────────
def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append(f"[슬라이드 {i}]\n" + "\n".join(slide_texts))
    return clean_extracted_text("\n\n".join(texts))


# ─── HWPX (XML 기반 한글 파일) ─────────────────────────
def _extract_hwpx(content: bytes) -> str:
    texts = []
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            # HWPX 내부에서 본문 XML 파일 찾기
            for name in sorted(zf.namelist()):
                if name.startswith("Contents/") and name.endswith(".xml"):
                    xml_data = zf.read(name)
                    root = ElementTree.fromstring(xml_data)
                    # 모든 텍스트 노드 수집
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            texts.append(elem.text.strip())
    except Exception as e:
        logger.warning(f"HWPX 추출 중 오류: {e}")
        raise
    return clean_extracted_text("\n".join(texts))


# ─── CSV ───────────────────────────────────────────────
def _extract_csv(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    lines = []
    for row in reader:
        cells = [c.strip() for c in row]
        if any(cells):
            lines.append(" | ".join(cells))
    return clean_extracted_text("\n".join(lines))


# ─── JSON ──────────────────────────────────────────────
def _extract_json(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    try:
        data = json.loads(text)
        return clean_extracted_text(json.dumps(data, ensure_ascii=False, indent=2))
    except json.JSONDecodeError:
        return clean_extracted_text(text)


# ─── XML / HTML ────────────────────────────────────────
def _extract_xml_html(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    import re
    # HTML 태그 제거
    clean = re.sub(r"<[^>]+>", " ", text)
    clean = re.sub(r"\s+", " ", clean)
    return clean_extracted_text(clean)


# ─── 플레인 텍스트 ────────────────────────────────────
def _extract_plain_text(content: bytes) -> str:
    # UTF-8 우선 시도, 실패 시 cp949 (한국어 레거시)
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            text = content.decode(enc)
            return clean_extracted_text(text)
        except (UnicodeDecodeError, ValueError):
            continue
    return clean_extracted_text(content.decode("utf-8", errors="replace"))


# ─── 유틸 ─────────────────────────────────────────────
def _ext(filename: str) -> str:
    return Path(filename).suffix.lower()
